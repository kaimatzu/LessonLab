import csv
import pandas as pd
import fitz
import requests
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

def scrape_txt(path: str):
    text_data = ""
    with open(path, 'r', encoding='utf-8') as file:
        text_data = file.read()
    return text_data

def scrape_pdf(path: str):
    pdf_document = fitz.open(path)
    text_data = ""
    for page_number in range(pdf_document.page_count):
        page = pdf_document[page_number]
        page_text = page.get_text()
        text_data += page_text
    return text_data

def scrape_docx(path: str):
    doc = Document(path)
    text_data = ""
    for paragraph in doc.paragraphs:
        text_data += paragraph.text
    return text_data

def scrape_pptx(path: str):
    ppt = Presentation(path)
    text_data = ""
    for slide_number, slide in enumerate(ppt.slides):
        text_data += f"\nSlide {slide_number+1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_data += shape.text + '\n'
    return text_data

def scrape_csv(path: str):
    text_data = ""
    with open(path, 'r') as csv_file:
        reader = csv.reader(csv_file)
        for row in reader:
            text_data += ','.join(row)
    return text_data

def scrape_xlsx(path: str):
    df = pd.read_excel(path)
    text_data = df.to_string(index = False)
    return text_data

def scrape_url(url: str):
    text_data = ""
    if validate_url(url) == False:
        webpage = requests.get(url)
        soup = BeautifulSoup(webpage.text, "html.parser")

        def contains_link(element):
            return element.find('a') is not None

        def is_valid_element(element):
            if element.name == 'a':
                return False
            if element.find_parent('a'):
                return False
            if element.find_parent('nav'):
                return False
            if element.find_parent('div', id=lambda value: value and ('menu' in value or 'footer' in value)):
                return False
            return True

        valid_elements = filter(is_valid_element, soup.recursiveChildGenerator())
        valid = ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'pre', 'li', 'span', 'div']

        for element in valid_elements:
            if element.name and element.name in valid and not contains_link(element):
                text_data += element.text.strip() + "\n"
    else:
        text_data = "Website is protected!"

    return text_data

def validate_url(url: str):
    # Add security keywords here
    security_keywords = ["cloudflare", "captcha"]

    # Check if scraped text is a response from security checks
    url = url.lower()
    for keyword in security_keywords:
        if keyword in url:
            return False
    if not url.strip():
        return False

    # Return true if website is scrapable and ready to feed to GPT-4
    return True